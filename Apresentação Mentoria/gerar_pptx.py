"""
Gera a apresentação Mentoria Plano FT como arquivo .pptx
Fernanda Trindade — Fisioterapeuta Pediátrica & Mentora de Negócios
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.oxml.ns as nsmap
from lxml import etree
import copy

# ── Cores da marca ────────────────────────────────────────────
NAVY       = RGBColor(0x1D, 0x2F, 0x55)
BROWN_DEEP = RGBColor(0x9A, 0x50, 0x2F)
BROWN_MID  = RGBColor(0xC2, 0x90, 0x81)
BROWN_LITE = RGBColor(0xBD, 0x7F, 0x68)
GREEN_PALE = RGBColor(0xD3, 0xD8, 0xAD)
GREEN_OLIVE= RGBColor(0x9D, 0x9D, 0x70)
GREEN_MID  = RGBColor(0x95, 0xAF, 0x71)
TEAL       = RGBColor(0x1F, 0x57, 0x7B)
CHARCOAL   = RGBColor(0x54, 0x54, 0x54)
CREAM      = RGBColor(0xF8, 0xF5, 0xF0)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)

# Slide 16:9
W = Inches(13.33)
H = Inches(7.5)

# Fonts
SERIF = "Cormorant Garamond"
SANS  = "Century Gothic"


def new_prs():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    return prs


def blank_layout(prs):
    return prs.slide_layouts[6]  # Blank


def rgb(color: RGBColor):
    return color


# ── Helpers ───────────────────────────────────────────────────

def fill_bg(slide, color: RGBColor):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color=None,
             line_color=None, line_width=Pt(0)):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.line.width = line_width
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, left, top, width, height,
             font_name=SANS, font_size=Pt(12), bold=False, italic=False,
             color=WHITE, align=PP_ALIGN.LEFT, word_wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = font_size
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_label(slide, text, left, top, width, color=BROWN_MID):
    return add_text(slide, text.upper(), left, top, width, Inches(0.35),
                    font_name=SANS, font_size=Pt(9), bold=True,
                    color=color, align=PP_ALIGN.LEFT)


def add_divider(slide, left, top, width=Inches(0.6), color=GREEN_MID):
    return add_rect(slide, left, top, width, Pt(1.5), fill_color=color)


def add_bullet_list(slide, items, left, top, width, height,
                    font_size=Pt(12), color=WHITE, prefix="✦  ",
                    spacing_pt=6, font_name=SANS):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(spacing_pt)
        run = p.add_run()
        run.text = prefix + item
        run.font.name = font_name
        run.font.size = font_size
        run.font.color.rgb = color
    return txBox


# ── Slide 1 — COVER ──────────────────────────────────────────
def s1_cover(prs):
    slide = prs.slides.add_slide(blank_layout(prs))
    fill_bg(slide, NAVY)

    # Left accent bar
    add_rect(slide, 0, 0, Inches(0.08), H, fill_color=GREEN_MID)

    # MENTORIA label
    add_text(slide, "MENTORIA", Inches(1.5), Inches(1.4), Inches(10), Inches(0.6),
             font_name=SANS, font_size=Pt(13), bold=True,
             color=GREEN_PALE, align=PP_ALIGN.LEFT)

    # PLANO FT — giant
    txBox = slide.shapes.add_textbox(Inches(1.5), Inches(1.9), Inches(10), Inches(2.2))
    tf = txBox.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "PLANO FT"
    run.font.name = SERIF
    run.font.size = Pt(108)
    run.font.bold = False
    run.font.color.rgb = WHITE

    # Tagline
    add_text(slide,
             "De fisioterapeuta invisível a referência lucrativa no particular",
             Inches(1.5), Inches(4.3), Inches(8), Inches(0.7),
             font_name=SANS, font_size=Pt(14), italic=True,
             color=RGBColor(0xFF,0xFF,0xFF), align=PP_ALIGN.LEFT)

    # Divider
    add_rect(slide, Inches(1.5), Inches(5.1), Inches(0.7), Pt(1.5), fill_color=GREEN_MID)

    # Author
    add_text(slide, "Fernanda Trindade  ·  Fisioterapeuta Pediátrica & Mentora de Negócios",
             Inches(1.5), Inches(5.4), Inches(9), Inches(0.5),
             font_name=SANS, font_size=Pt(10),
             color=RGBColor(0xAA, 0xBB, 0xCC), align=PP_ALIGN.LEFT)

    # Right botanical accent rect (decorative)
    add_rect(slide, Inches(11.5), 0, Inches(1.83), H,
             fill_color=RGBColor(0x1A, 0x29, 0x4C))


# ── Slide 2 — SOBRE FERNANDA ─────────────────────────────────
def s2_sobre(prs):
    slide = prs.slides.add_slide(blank_layout(prs))
    fill_bg(slide, CREAM)

    # Right navy panel
    add_rect(slide, Inches(8.5), 0, Inches(4.83), H, fill_color=NAVY)

    # Left content
    add_text(slide, "QUEM É", Inches(0.8), Inches(1.2), Inches(5), Inches(0.4),
             font_name=SANS, font_size=Pt(9), bold=True,
             color=BROWN_DEEP, align=PP_ALIGN.LEFT)

    add_text(slide, "Fernanda Trindade",
             Inches(0.8), Inches(1.7), Inches(7), Inches(1.2),
             font_name=SERIF, font_size=Pt(52), color=NAVY, align=PP_ALIGN.LEFT)

    add_text(slide, "Fisioterapeuta Pediátrica & Mentora de Negócios",
             Inches(0.8), Inches(2.9), Inches(7), Inches(0.5),
             font_name=SANS, font_size=Pt(12), color=GREEN_OLIVE, align=PP_ALIGN.LEFT)

    add_divider(slide, Inches(0.8), Inches(3.55), color=BROWN_MID)

    bio = ("Fisioterapeuta pediátrica com anos de experiência clínica, "
           "Fernanda trilhou o caminho que suas mentoradas agora percorrem — "
           "do trabalho excessivo e baixo lucro para um negócio estruturado, "
           "previsível e lucrativo. Hoje, além de atender, ela guia "
           "fisioterapeutas pediátricas a construírem negócios que refletem "
           "quem elas são e o quanto elas valem.")
    add_text(slide, bio,
             Inches(0.8), Inches(3.8), Inches(7.2), Inches(2.5),
             font_name=SANS, font_size=Pt(12), color=CHARCOAL, align=PP_ALIGN.LEFT)

    # Right panel — photo placeholder
    add_text(slide, "[ foto ]",
             Inches(8.8), Inches(3.3), Inches(4), Inches(0.5),
             font_name=SANS, font_size=Pt(11),
             color=RGBColor(0x88,0x99,0xAA), align=PP_ALIGN.CENTER)

    add_text(slide, "Fernanda Trindade",
             Inches(8.8), Inches(6.4), Inches(4), Inches(0.5),
             font_name=SANS, font_size=Pt(9),
             color=RGBColor(0xAA,0xBB,0xCC), align=PP_ALIGN.CENTER)


# ── Slide 3 — COMO SURGIU ────────────────────────────────────
def s3_surgiu(prs):
    slide = prs.slides.add_slide(blank_layout(prs))
    fill_bg(slide, GREEN_MID)

    add_text(slide, "A ORIGEM DO MÉTODO",
             Inches(0.5), Inches(0.6), Inches(12), Inches(0.4),
             font_name=SANS, font_size=Pt(9),
             color=RGBColor(0xFF,0xFF,0xFF), align=PP_ALIGN.CENTER)

    add_text(slide, "Como nasceu o Plano FT",
             Inches(0.5), Inches(1.0), Inches(12), Inches(1.1),
             font_name=SERIF, font_size=Pt(52), italic=True,
             color=WHITE, align=PP_ALIGN.CENTER)

    cards = [
        ("1", "Eu vivi isso", "Fernanda passou pelo mesmo ciclo de exaustão e pouco retorno financeiro"),
        ("2", "Aprendi na prática", "Estruturou seu próprio negócio do zero, testando e validando cada etapa"),
        ("3", "Agora compartilho", "Transformou sua experiência em um método replicável e personalizado"),
    ]

    left_positions = [Inches(0.7), Inches(4.7), Inches(8.7)]
    for (num, title, text), lpos in zip(cards, left_positions):
        # Card background
        add_rect(slide, lpos, Inches(2.5), Inches(3.8), Inches(4.0),
                 fill_color=RGBColor(0x7A, 0x9A, 0x5A))

        add_text(slide, num, lpos + Inches(0.2), Inches(2.65), Inches(1), Inches(1.0),
                 font_name=SERIF, font_size=Pt(56),
                 color=RGBColor(0xFF,0xFF,0xFF), align=PP_ALIGN.LEFT)

        add_text(slide, title, lpos + Inches(0.2), Inches(3.65), Inches(3.4), Inches(0.5),
                 font_name=SERIF, font_size=Pt(20), bold=False,
                 color=WHITE, align=PP_ALIGN.LEFT)

        add_text(slide, text, lpos + Inches(0.2), Inches(4.25), Inches(3.4), Inches(1.8),
                 font_name=SANS, font_size=Pt(11),
                 color=RGBColor(0xF0,0xF8,0xF0), align=PP_ALIGN.LEFT)


# ── Slide 4 — A PROMESSA ─────────────────────────────────────
def s4_promessa(prs):
    slide = prs.slides.add_slide(blank_layout(prs))
    fill_bg(slide, CREAM)

    add_text(slide, "A PROMESSA",
             Inches(0.5), Inches(0.8), Inches(12), Inches(0.4),
             font_name=SANS, font_size=Pt(9),
             color=BROWN_MID, align=PP_ALIGN.CENTER)

    quote = "\"Você não precisa atender mais...\nvocê precisa estruturar o que já tem.\""
    add_text(slide, quote,
             Inches(1.0), Inches(1.4), Inches(11), Inches(2.4),
             font_name=SERIF, font_size=Pt(40), italic=True,
             color=NAVY, align=PP_ALIGN.CENTER)

    add_divider(slide, Inches(6.2), Inches(3.9), color=GREEN_OLIVE)

    sub = ("Em 4 meses, você vai sair de fisioterapeuta operacional para profissional "
           "com negócio estruturado, lucrativo e previsível.")
    add_text(slide, sub,
             Inches(1.5), Inches(4.2), Inches(10), Inches(1.5),
             font_name=SANS, font_size=Pt(14),
             color=CHARCOAL, align=PP_ALIGN.CENTER)


# ── Slide 5 — PARA QUEM É ────────────────────────────────────
def s5_paraquem(prs):
    slide = prs.slides.add_slide(blank_layout(prs))

    # Left panel (navy)
    add_rect(slide, 0, 0, Inches(7.8), H, fill_color=NAVY)
    # Right panel (cream)
    add_rect(slide, Inches(7.8), 0, Inches(5.53), H, fill_color=CREAM)

    add_text(slide, "ESTA MENTORIA É PARA VOCÊ SE...",
             Inches(0.6), Inches(0.8), Inches(7), Inches(0.4),
             font_name=SANS, font_size=Pt(9),
             color=GREEN_PALE, align=PP_ALIGN.LEFT)

    bullets = [
        "Já atua como fisioterapeuta pediátrica mas trabalha muito e ganha pouco",
        "Tem dificuldade de precificar seus serviços com segurança",
        "Quer se posicionar como referência na sua cidade",
        "Está cansada de atender em várias cidades sem previsibilidade",
        "Quer uma agenda estratégica e um negócio que sustente sua qualidade de vida",
    ]
    add_bullet_list(slide, bullets,
                    Inches(0.6), Inches(1.4), Inches(7), Inches(5.5),
                    font_size=Pt(13.5), color=RGBColor(0xEE,0xEE,0xEE),
                    prefix="✦  ", spacing_pt=10)

    # Right panel quote
    add_text(slide, "\"A mentoria certa,\nno momento certo,\nmuda tudo.\"",
             Inches(8.1), Inches(2.2), Inches(5), Inches(3),
             font_name=SERIF, font_size=Pt(28), italic=True,
             color=NAVY, align=PP_ALIGN.CENTER)


# ── Slide 6 — VOCÊ SE RECONHECE? ─────────────────────────────
def s6_reconhece(prs):
    slide = prs.slides.add_slide(blank_layout(prs))
    fill_bg(slide, BROWN_LITE)

    add_text(slide, "VOCÊ JÁ DISSE ALGUMA DESSAS FRASES?",
             Inches(0.5), Inches(0.55), Inches(12), Inches(0.5),
             font_name=SANS, font_size=Pt(9),
             color=RGBColor(0xFF,0xFF,0xFF), align=PP_ALIGN.CENTER)

    add_text(slide, "Você se reconhece?",
             Inches(0.5), Inches(1.0), Inches(12), Inches(0.9),
             font_name=SERIF, font_size=Pt(44), italic=True,
             color=WHITE, align=PP_ALIGN.CENTER)

    quotes = [
        '"Eu trabalho muito e não vejo resultado"',
        '"Eu sei atender, mas não sei crescer"',
        '"Eu não consigo cobrar mais"',
        '"Não tenho organização nenhuma"',
        '"Eu vivo cansada"',
    ]
    positions = [
        (Inches(0.5), Inches(2.2)),
        (Inches(4.8), Inches(2.2)),
        (Inches(9.1), Inches(2.2)),
        (Inches(1.8), Inches(4.5)),
        (Inches(7.0), Inches(4.5)),
    ]
    for (txt, (lpos, tpos)) in zip(quotes, positions):
        add_rect(slide, lpos, tpos, Inches(3.8), Inches(1.85),
                 fill_color=RGBColor(0xD0, 0x8F, 0x78))
        add_text(slide, txt, lpos + Inches(0.25), tpos + Inches(0.3),
                 Inches(3.3), Inches(1.3),
                 font_name=SERIF, font_size=Pt(15), italic=True,
                 color=WHITE, align=PP_ALIGN.CENTER)


# ── Slide 7 — HOJE vs DEPOIS ─────────────────────────────────
def s7_hoje_depois(prs):
    slide = prs.slides.add_slide(blank_layout(prs))

    add_rect(slide, 0, 0, Inches(6.2), H, fill_color=NAVY)
    add_rect(slide, Inches(6.2), 0, Inches(0.6), H, fill_color=WHITE)
    add_rect(slide, Inches(6.8), 0, Inches(6.53), H, fill_color=GREEN_MID)

    # Center arrow
    add_text(slide, "→", Inches(6.1), Inches(3.4), Inches(0.8), Inches(0.5),
             font_name=SANS, font_size=Pt(22), color=NAVY, align=PP_ALIGN.CENTER)

    # HOJE side
    add_text(slide, "HOJE", Inches(0.5), Inches(0.6), Inches(5), Inches(0.4),
             font_name=SANS, font_size=Pt(9), color=RGBColor(0xAA,0xBB,0xCC), align=PP_ALIGN.LEFT)
    add_text(slide, "Hoje", Inches(0.5), Inches(0.9), Inches(5), Inches(0.9),
             font_name=SERIF, font_size=Pt(54), color=WHITE, align=PP_ALIGN.LEFT)

    hoje_items = [
        "Trabalha em excesso",
        "Atende em várias cidades",
        "Cobra abaixo do valor",
        "Agenda desorganizada",
        "Sem posicionamento claro",
        "Financeiro caótico",
    ]
    add_bullet_list(slide, hoje_items, Inches(0.5), Inches(2.0), Inches(5.5), Inches(4.5),
                    font_size=Pt(13.5), color=RGBColor(0xDD,0xDD,0xDD),
                    prefix="—  ", spacing_pt=8)

    # DEPOIS side
    add_text(slide, "DEPOIS", Inches(7.2), Inches(0.6), Inches(5), Inches(0.4),
             font_name=SANS, font_size=Pt(9), color=RGBColor(0xFF,0xFF,0xFF), align=PP_ALIGN.LEFT)
    add_text(slide, "Depois", Inches(7.2), Inches(0.9), Inches(5), Inches(0.9),
             font_name=SERIF, font_size=Pt(54), color=WHITE, align=PP_ALIGN.LEFT)

    depois_items = [
        "Trabalha menos com mais retorno",
        "Agenda concentrada e estratégica",
        "Cobra premium com segurança",
        "Agenda previsível e organizada",
        "Referência reconhecida na cidade",
        "Lucro real e estruturado",
    ]
    add_bullet_list(slide, depois_items, Inches(7.2), Inches(2.0), Inches(5.5), Inches(4.5),
                    font_size=Pt(13.5), color=WHITE,
                    prefix="✓  ", spacing_pt=8)


# ── Slide 8 — ESTRUTURA GERAL ─────────────────────────────────
def s8_estrutura(prs):
    slide = prs.slides.add_slide(blank_layout(prs))
    fill_bg(slide, CREAM)

    add_text(slide, "A JORNADA", Inches(0.5), Inches(0.55), Inches(12), Inches(0.4),
             font_name=SANS, font_size=Pt(9), color=BROWN_MID, align=PP_ALIGN.CENTER)
    add_text(slide, "A Estrutura da Mentoria",
             Inches(0.5), Inches(0.9), Inches(12), Inches(0.9),
             font_name=SERIF, font_size=Pt(44), color=NAVY, align=PP_ALIGN.CENTER)

    nums = [("4", "Meses"), ("12", "Encontros"), ("5", "Pilares")]
    positions = [Inches(1.5), Inches(5.2), Inches(8.9)]
    for (num, label), lpos in zip(nums, positions):
        add_text(slide, num, lpos, Inches(2.0), Inches(3.0), Inches(2.2),
                 font_name=SERIF, font_size=Pt(120),
                 color=BROWN_DEEP, align=PP_ALIGN.CENTER)
        add_text(slide, label.upper(), lpos, Inches(4.3), Inches(3.0), Inches(0.5),
                 font_name=SANS, font_size=Pt(11), bold=True,
                 color=NAVY, align=PP_ALIGN.CENTER)

    # Dividers between numbers
    for x in [Inches(4.7), Inches(8.4)]:
        add_rect(slide, x, Inches(2.2), Pt(1), Inches(2.5),
                 fill_color=RGBColor(0xCC,0xCC,0xCC))

    add_text(slide,
             "Uma jornada estratégica e progressiva, do alicerce ao próximo nível.",
             Inches(1.0), Inches(5.1), Inches(11), Inches(0.6),
             font_name=SANS, font_size=Pt(13),
             color=CHARCOAL, align=PP_ALIGN.CENTER)


# ── Slide 9 — 5 PILARES OVERVIEW ─────────────────────────────
def s9_pilares(prs):
    slide = prs.slides.add_slide(blank_layout(prs))
    fill_bg(slide, NAVY)

    add_text(slide, "OS 5 PILARES", Inches(0.5), Inches(0.5), Inches(12), Inches(0.4),
             font_name=SANS, font_size=Pt(9), color=GREEN_PALE, align=PP_ALIGN.CENTER)
    add_text(slide, "A Estrutura Completa",
             Inches(0.5), Inches(0.85), Inches(12), Inches(0.9),
             font_name=SERIF, font_size=Pt(44), color=WHITE, align=PP_ALIGN.CENTER)

    pilares = [
        (GREEN_MID,  "1", "Clareza e Identidade\nEstratégica",      "2 encontros · Mês 1"),
        (BROWN_DEEP, "2", "Precificação e\nEstrutura de Negócio",   "2 encontros · Mês 1-2"),
        (TEAL,       "3", "Presença Digital e\nAtração de Pacientes","3 encontros · Mês 2-3"),
        (BROWN_MID,  "4", "Organização e\nControle do Negócio",     "2 encontros · Mês 3"),
        (GREEN_OLIVE,"5", "Consolidação e\nPróximo Nível",           "3 encontros · Mês 4"),
    ]

    card_w = Inches(2.35)
    gap = Inches(0.18)
    start = Inches(0.5)
    for i, (color, num, title, meta) in enumerate(pilares):
        lpos = start + i * (card_w + gap)
        add_rect(slide, lpos, Inches(2.0), card_w, Inches(4.8), fill_color=color)
        add_text(slide, num, lpos + Inches(0.18), Inches(2.15), Inches(1.5), Inches(1.1),
                 font_name=SERIF, font_size=Pt(64),
                 color=RGBColor(0xFF,0xFF,0xFF), align=PP_ALIGN.LEFT)
        add_text(slide, title, lpos + Inches(0.18), Inches(3.3), card_w - Inches(0.36), Inches(1.6),
                 font_name=SERIF, font_size=Pt(15.5),
                 color=WHITE, align=PP_ALIGN.LEFT)
        add_text(slide, meta, lpos + Inches(0.18), Inches(5.9), card_w - Inches(0.36), Inches(0.6),
                 font_name=SANS, font_size=Pt(9.5),
                 color=RGBColor(0xFF,0xFF,0xFF), align=PP_ALIGN.LEFT)


# ── Slides 10–14 — PILAR DETAIL helper ────────────────────────
def pilar_slide(prs, accent_color, num, title, sub, month_info,
                enc1_num, enc1_title, enc1_items,
                enc2_num=None, enc2_title=None, enc2_items=None,
                enc3_num=None, enc3_title=None, enc3_items=None,
                highlight=None):
    slide = prs.slides.add_slide(blank_layout(prs))
    fill_bg(slide, CREAM)

    # Left accent bar
    add_rect(slide, 0, 0, Inches(0.1), H, fill_color=accent_color)

    # Big number
    add_text(slide, num, Inches(0.4), Inches(0.35), Inches(1.8), Inches(1.5),
             font_name=SERIF, font_size=Pt(96),
             color=accent_color, align=PP_ALIGN.LEFT)

    # Title block
    add_text(slide, "PILAR " + num, Inches(2.1), Inches(0.6), Inches(10), Inches(0.35),
             font_name=SANS, font_size=Pt(9), bold=True,
             color=RGBColor(0x99,0x99,0x99), align=PP_ALIGN.LEFT)
    add_text(slide, title, Inches(2.1), Inches(0.85), Inches(9), Inches(0.9),
             font_name=SERIF, font_size=Pt(32),
             color=NAVY, align=PP_ALIGN.LEFT)
    add_text(slide, sub, Inches(2.1), Inches(1.75), Inches(9), Inches(0.4),
             font_name=SANS, font_size=Pt(11), italic=True,
             color=CHARCOAL, align=PP_ALIGN.LEFT)
    add_text(slide, month_info, Inches(2.1), Inches(2.15), Inches(9), Inches(0.35),
             font_name=SANS, font_size=Pt(9),
             color=accent_color, align=PP_ALIGN.LEFT)

    # Encontro cards
    enc_data = [(enc1_num, enc1_title, enc1_items)]
    if enc2_num:
        enc_data.append((enc2_num, enc2_title, enc2_items))
    if enc3_num:
        enc_data.append((enc3_num, enc3_title, enc3_items))

    card_w = (Inches(12.7) - Inches(0.5) * (len(enc_data) - 1)) / len(enc_data)
    start_l = Inches(0.3)

    for i, (enum, etitle, eitems) in enumerate(enc_data):
        lpos = start_l + i * (card_w + Inches(0.25))
        add_rect(slide, lpos, Inches(2.7), card_w, Inches(4.4),
                 fill_color=WHITE, line_color=RGBColor(0xE8,0xE4,0xDF), line_width=Pt(0.5))
        # Top border strip
        add_rect(slide, lpos, Inches(2.7), card_w, Pt(3.5), fill_color=accent_color)

        add_text(slide, enum, lpos + Inches(0.2), Inches(2.85), card_w, Inches(0.3),
                 font_name=SANS, font_size=Pt(8.5), bold=True,
                 color=accent_color, align=PP_ALIGN.LEFT)
        add_text(slide, etitle, lpos + Inches(0.2), Inches(3.1), card_w - Inches(0.4), Inches(0.6),
                 font_name=SERIF, font_size=Pt(17),
                 color=NAVY, align=PP_ALIGN.LEFT)

        if eitems:
            add_bullet_list(slide, eitems,
                            lpos + Inches(0.2), Inches(3.75),
                            card_w - Inches(0.4), Inches(2.8),
                            font_size=Pt(10.5), color=CHARCOAL,
                            prefix="·  ", spacing_pt=4)

    if highlight:
        add_rect(slide, Inches(0.3), Inches(7.1), Inches(12.7), Inches(0.32),
                 fill_color=NAVY)
        add_text(slide, highlight, Inches(0.5), Inches(7.1), Inches(12.3), Inches(0.32),
                 font_name=SANS, font_size=Pt(10), italic=True,
                 color=GREEN_PALE, align=PP_ALIGN.LEFT)

    return slide


# ── Slide 15 — TIMELINE ──────────────────────────────────────
def s15_timeline(prs):
    slide = prs.slides.add_slide(blank_layout(prs))
    fill_bg(slide, NAVY)

    add_text(slide, "CRONOGRAMA", Inches(0.5), Inches(0.5), Inches(12), Inches(0.4),
             font_name=SANS, font_size=Pt(9), color=RGBColor(0x88,0x99,0xAA), align=PP_ALIGN.CENTER)
    add_text(slide, "A Jornada Mês a Mês",
             Inches(0.5), Inches(0.85), Inches(12), Inches(0.9),
             font_name=SERIF, font_size=Pt(44), color=WHITE, align=PP_ALIGN.CENTER)

    # Timeline line
    add_rect(slide, Inches(0.8), Inches(2.5), Inches(11.7), Pt(1.5),
             fill_color=RGBColor(0x44, 0x55, 0x77))

    months = [
        (GREEN_MID,  "Mês 1", "O Alicerce",
         "Pilares 1 e 2\nEncontros 1, 2, 3, 4\nClareza, identidade, precificação"),
        (BROWN_DEEP, "Mês 2", "A Vitrine",
         "Pilar 3 início\nEncontros 5, 6, 7\nInstagram, conteúdo, conversão"),
        (TEAL,       "Mês 3", "A Estrutura",
         "Pilares 3 fim e 4\nEncontros 8, 9\nOrganização, gestão, indicadores"),
        (BROWN_MID,  "Mês 4", "O Futuro",
         "Pilar 5\nEncontros 10, 11, 12\nAjustes, expansão, celebração"),
    ]

    dot_w = Inches(0.7)
    cols = [Inches(1.0), Inches(4.2), Inches(7.4), Inches(10.6)]

    for (color, mes, title, detail), lpos in zip(months, cols):
        # Dot
        add_rect(slide, lpos + Inches(1.0), Inches(2.15), dot_w, dot_w,
                 fill_color=color)
        add_text(slide, "●", lpos + Inches(0.95), Inches(2.1), Inches(0.9), Inches(0.6),
                 font_name=SANS, font_size=Pt(32), color=color, align=PP_ALIGN.CENTER)

        add_text(slide, mes, lpos, Inches(2.9), Inches(3.0), Inches(0.4),
                 font_name=SANS, font_size=Pt(9.5), bold=True,
                 color=color, align=PP_ALIGN.LEFT)
        add_text(slide, title, lpos, Inches(3.25), Inches(3.0), Inches(0.6),
                 font_name=SERIF, font_size=Pt(22), color=WHITE, align=PP_ALIGN.LEFT)

        add_rect(slide, lpos, Inches(4.0), Inches(3.0), Inches(3.0),
                 fill_color=RGBColor(0x25, 0x3B, 0x65))
        add_text(slide, detail, lpos + Inches(0.15), Inches(4.1), Inches(2.7), Inches(2.8),
                 font_name=SANS, font_size=Pt(10.5),
                 color=RGBColor(0xCC,0xD8,0xEE), align=PP_ALIGN.LEFT)


# ── Slide 16 — ENTREGÁVEIS ────────────────────────────────────
def s16_entregaveis(prs):
    slide = prs.slides.add_slide(blank_layout(prs))
    fill_bg(slide, CREAM)

    add_text(slide, "O QUE VOCÊ RECEBE", Inches(0.5), Inches(0.55), Inches(12), Inches(0.4),
             font_name=SANS, font_size=Pt(9), color=BROWN_MID, align=PP_ALIGN.CENTER)
    add_text(slide, "O que está incluído",
             Inches(0.5), Inches(0.9), Inches(12), Inches(0.8),
             font_name=SERIF, font_size=Pt(40), color=NAVY, align=PP_ALIGN.CENTER)

    items = [
        ("12 encontros estratégicos ao vivo", GREEN_MID),
        ("Materiais exclusivos por pilar", BROWN_DEEP),
        ("Suporte entre os encontros", TEAL),
        ("Exercícios práticos com aplicação imediata", BROWN_MID),
        ("Calendário de conteúdo personalizado", GREEN_OLIVE),
        ("Scripts e roteiros de atendimento", BROWN_LITE),
        ("Ferramentas de gestão financeira", GREEN_MID),
        ("Plano de continuidade pós-mentoria", NAVY),
        ("Acesso a comunidade de fisioterapeutas", BROWN_DEEP),
        ("Gravação de todos os encontros", TEAL),
    ]

    cols = 5
    card_w = Inches(2.4)
    card_h = Inches(1.55)
    gap_h = Inches(0.18)
    gap_v = Inches(0.15)
    start_l = Inches(0.35)
    start_t = Inches(1.85)

    for i, (text, color) in enumerate(items):
        row = i // cols
        col = i % cols
        lpos = start_l + col * (card_w + gap_h)
        tpos = start_t + row * (card_h + gap_v)
        add_rect(slide, lpos, tpos, card_w, card_h, fill_color=WHITE)
        # Bottom accent
        add_rect(slide, lpos, tpos + card_h - Pt(3), card_w, Pt(3), fill_color=color)
        add_text(slide, "✓", lpos + Inches(0.15), tpos + Inches(0.12), Inches(0.4), Inches(0.4),
                 font_name=SANS, font_size=Pt(14), color=color, align=PP_ALIGN.CENTER)
        add_text(slide, text, lpos + Inches(0.12), tpos + Inches(0.52),
                 card_w - Inches(0.24), Inches(0.9),
                 font_name=SANS, font_size=Pt(10), color=NAVY, align=PP_ALIGN.LEFT)


# ── Slide 17 — TRANSFORMAÇÃO ─────────────────────────────────
def s17_transformacao(prs):
    slide = prs.slides.add_slide(blank_layout(prs))
    fill_bg(slide, GREEN_OLIVE)

    add_text(slide, "A JORNADA COMPLETA", Inches(0.5), Inches(0.55), Inches(12), Inches(0.4),
             font_name=SANS, font_size=Pt(9),
             color=RGBColor(0xFF,0xFF,0xFF), align=PP_ALIGN.CENTER)
    add_text(slide, "A transformação que acontece",
             Inches(0.5), Inches(0.9), Inches(12), Inches(0.9),
             font_name=SERIF, font_size=Pt(40), italic=True,
             color=WHITE, align=PP_ALIGN.CENTER)

    antes_items = [
        "Exaustão — trabalha muito, sobra pouco",
        "Medo de cobrar mais e perder pacientes",
        "Posta sem estratégia, sem retorno",
        "Agenda desorganizada e imprevisível",
        "Sem controle financeiro real",
        "Invisível na sua cidade",
    ]
    depois_items = [
        "Clareza total de quem é e o que oferece",
        "Cobra premium com segurança e naturalidade",
        "Conteúdo estratégico que atrai e converte",
        "Agenda previsível e semana ideal estruturada",
        "Lucro real e controle financeiro funcionando",
        "Reconhecida como referência na cidade",
    ]

    add_rect(slide, Inches(0.3), Inches(2.0), Inches(5.8), Inches(4.5),
             fill_color=RGBColor(0x88, 0x8D, 0x60))
    add_text(slide, "Antes da Mentoria",
             Inches(0.5), Inches(2.1), Inches(5), Inches(0.4),
             font_name=SANS, font_size=Pt(9), color=RGBColor(0xEE,0xFF,0xEE), align=PP_ALIGN.LEFT)
    add_bullet_list(slide, antes_items, Inches(0.5), Inches(2.55), Inches(5.6), Inches(3.7),
                    font_size=Pt(12), color=RGBColor(0xF0,0xF8,0xF0),
                    prefix="✕  ", spacing_pt=7)

    add_text(slide, "→", Inches(6.3), Inches(3.8), Inches(0.7), Inches(0.7),
             font_name=SANS, font_size=Pt(28), color=WHITE, align=PP_ALIGN.CENTER)

    add_rect(slide, Inches(7.2), Inches(2.0), Inches(5.8), Inches(4.5),
             fill_color=RGBColor(0x88, 0x8D, 0x60))
    add_text(slide, "Depois da Mentoria",
             Inches(7.4), Inches(2.1), Inches(5), Inches(0.4),
             font_name=SANS, font_size=Pt(9), color=RGBColor(0xEE,0xFF,0xEE), align=PP_ALIGN.LEFT)
    add_bullet_list(slide, depois_items, Inches(7.4), Inches(2.55), Inches(5.6), Inches(3.7),
                    font_size=Pt(12), color=WHITE,
                    prefix="✓  ", spacing_pt=7)

    add_text(slide,
             '"Você não compra conteúdo. Você compra direção, estrutura e segurança para crescer."',
             Inches(1.0), Inches(6.65), Inches(11), Inches(0.6),
             font_name=SERIF, font_size=Pt(14), italic=True,
             color=RGBColor(0xF0,0xF8,0xF0), align=PP_ALIGN.CENTER)


# ── Slide 18 — INVESTIMENTO ──────────────────────────────────
def s18_investimento(prs):
    slide = prs.slides.add_slide(blank_layout(prs))
    fill_bg(slide, NAVY)

    add_text(slide, "INVESTIMENTO", Inches(0.5), Inches(0.55), Inches(12), Inches(0.4),
             font_name=SANS, font_size=Pt(9), color=GREEN_PALE, align=PP_ALIGN.CENTER)
    add_text(slide, "Mentoria Plano FT",
             Inches(0.5), Inches(0.9), Inches(12), Inches(0.9),
             font_name=SERIF, font_size=Pt(44), color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, "4 meses de transformação estratégica",
             Inches(0.5), Inches(1.85), Inches(12), Inches(0.45),
             font_name=SANS, font_size=Pt(12),
             color=RGBColor(0x88,0x99,0xAA), align=PP_ALIGN.CENTER)

    # Investment box
    add_rect(slide, Inches(3.8), Inches(2.5), Inches(5.7), Inches(3.6),
             fill_color=RGBColor(0x22, 0x35, 0x5F),
             line_color=RGBColor(0x3A, 0x50, 0x80), line_width=Pt(0.75))

    add_text(slide, "INVESTIMENTO TOTAL",
             Inches(4.0), Inches(2.7), Inches(5), Inches(0.35),
             font_name=SANS, font_size=Pt(9),
             color=RGBColor(0x88,0x99,0xAA), align=PP_ALIGN.CENTER)

    add_text(slide, "R$ 3.000",
             Inches(3.8), Inches(3.0), Inches(5.7), Inches(1.3),
             font_name=SERIF, font_size=Pt(80),
             color=GREEN_PALE, align=PP_ALIGN.CENTER)

    add_text(slide, "Ou em até 6x de R$ 500,00\nPagamento via PIX, cartão ou boleto",
             Inches(4.0), Inches(4.3), Inches(5.0), Inches(0.7),
             font_name=SANS, font_size=Pt(11.5),
             color=RGBColor(0xAA,0xBB,0xCC), align=PP_ALIGN.CENTER)

    add_rect(slide, Inches(5.5), Inches(5.2), Inches(2.3), Pt(1),
             fill_color=RGBColor(0x44,0x55,0x77))

    add_text(slide, "Vagas limitadas por turma para garantir acompanhamento personalizado",
             Inches(2.5), Inches(5.5), Inches(8.3), Inches(0.55),
             font_name=SANS, font_size=Pt(10.5),
             color=RGBColor(0x77,0x88,0x99), align=PP_ALIGN.CENTER)


# ── Slide 19 — FAQ ────────────────────────────────────────────
def s19_faq(prs):
    slide = prs.slides.add_slide(blank_layout(prs))
    fill_bg(slide, CREAM)

    add_text(slide, "DÚVIDAS FREQUENTES", Inches(0.5), Inches(0.55), Inches(12), Inches(0.4),
             font_name=SANS, font_size=Pt(9), color=BROWN_MID, align=PP_ALIGN.CENTER)
    add_text(slide, "Perguntas que você pode ter",
             Inches(0.5), Inches(0.9), Inches(12), Inches(0.8),
             font_name=SERIF, font_size=Pt(40), color=NAVY, align=PP_ALIGN.CENTER)

    faqs = [
        (GREEN_MID,  '"Agora não é o momento certo..."',
         'O momento nunca vai ser perfeito. Se você sente que algo precisa mudar, esse é exatamente o sinal. A mentoria foi desenhada para caber na sua realidade atual.'),
        (BROWN_DEEP, '"Acho que está caro para mim..."',
         'A pergunta certa não é "quanto custa", mas "quanto me custa continuar onde estou?". É um investimento direto no aumento do seu faturamento.'),
        (TEAL,       '"Minha realidade é diferente..."',
         'A mentoria é personalizada. O método é estruturado, mas cada encontro é adaptado à sua cidade, ao seu perfil e ao seu momento. Não existe solução genérica aqui.'),
        (BROWN_MID,  '"Tenho medo de não conseguir aplicar..."',
         'Por isso cada encontro tem tarefas práticas e imediatas. O método foi desenhado para ação, não para teoria. Você vai saber exatamente o que fazer.'),
    ]

    positions = [
        (Inches(0.3), Inches(2.0)),
        (Inches(6.8), Inches(2.0)),
        (Inches(0.3), Inches(5.0)),
        (Inches(6.8), Inches(5.0)),
    ]
    for (color, q, a), (lpos, tpos) in zip(faqs, positions):
        add_rect(slide, lpos, tpos, Inches(6.3), Inches(2.65),
                 fill_color=WHITE, line_color=color, line_width=Pt(3))
        # Left border strip
        add_rect(slide, lpos, tpos, Pt(4), Inches(2.65), fill_color=color)
        add_text(slide, q, lpos + Inches(0.2), tpos + Inches(0.15), Inches(5.9), Inches(0.55),
                 font_name=SERIF, font_size=Pt(14), italic=True, color=NAVY, align=PP_ALIGN.LEFT)
        add_text(slide, a, lpos + Inches(0.2), tpos + Inches(0.72), Inches(5.9), Inches(1.75),
                 font_name=SANS, font_size=Pt(11), color=CHARCOAL, align=PP_ALIGN.LEFT)


# ── Slide 20 — CTA FINAL ─────────────────────────────────────
def s20_cta(prs):
    slide = prs.slides.add_slide(blank_layout(prs))
    fill_bg(slide, NAVY)

    # Decorative right panel
    add_rect(slide, Inches(9.0), 0, Inches(4.33), H,
             fill_color=RGBColor(0x18, 0x28, 0x4A))
    add_rect(slide, Inches(9.0), 0, Pt(2), H, fill_color=GREEN_MID)

    add_text(slide, "PRONTA PARA DAR O PRÓXIMO PASSO?",
             Inches(0.5), Inches(0.7), Inches(8), Inches(0.4),
             font_name=SANS, font_size=Pt(9), color=GREEN_PALE, align=PP_ALIGN.LEFT)

    add_text(slide, '"Chegou a hora de estruturar\no que você já tem."',
             Inches(0.5), Inches(1.2), Inches(8.2), Inches(2.4),
             font_name=SERIF, font_size=Pt(44), italic=True,
             color=WHITE, align=PP_ALIGN.LEFT)

    add_text(slide, "De fisioterapeuta invisível a referência lucrativa no particular",
             Inches(0.5), Inches(3.75), Inches(8), Inches(0.6),
             font_name=SANS, font_size=Pt(13),
             color=RGBColor(0xAA,0xBB,0xCC), align=PP_ALIGN.LEFT)

    # CTA button (simulated)
    add_rect(slide, Inches(0.5), Inches(4.55), Inches(3.2), Inches(0.62),
             fill_color=BROWN_LITE)
    add_text(slide, "QUERO MINHA VAGA",
             Inches(0.5), Inches(4.57), Inches(3.2), Inches(0.55),
             font_name=SANS, font_size=Pt(11), bold=True,
             color=WHITE, align=PP_ALIGN.CENTER)

    add_rect(slide, Inches(0.5), Inches(5.4), Inches(2), Pt(1),
             fill_color=RGBColor(0x44,0x55,0x77))

    add_text(slide,
             "Fernanda Trindade  ·  Fisioterapeuta Pediátrica & Mentora\n"
             "@nandatrindade.fisio",
             Inches(0.5), Inches(5.6), Inches(7.5), Inches(0.85),
             font_name=SANS, font_size=Pt(11),
             color=RGBColor(0x77,0x88,0x99), align=PP_ALIGN.LEFT)

    add_text(slide, "Mentoria Plano FT",
             Inches(0.5), Inches(6.7), Inches(4), Inches(0.55),
             font_name=SERIF, font_size=Pt(22),
             color=RGBColor(0x55,0x66,0x88), align=PP_ALIGN.LEFT)

    # Right panel content
    add_text(slide, "FT",
             Inches(10.3), Inches(2.8), Inches(2.0), Inches(1.5),
             font_name=SERIF, font_size=Pt(96),
             color=RGBColor(0x2A, 0x3D, 0x6A), align=PP_ALIGN.CENTER)
    add_text(slide, "fisioterapeuta pediátrica",
             Inches(9.4), Inches(5.5), Inches(3.8), Inches(0.5),
             font_name=SANS, font_size=Pt(9),
             color=RGBColor(0x44,0x55,0x77), align=PP_ALIGN.CENTER)


# ── BUILD ALL SLIDES ──────────────────────────────────────────
def build():
    prs = new_prs()

    s1_cover(prs)
    s2_sobre(prs)
    s3_surgiu(prs)
    s4_promessa(prs)
    s5_paraquem(prs)
    s6_reconhece(prs)
    s7_hoje_depois(prs)
    s8_estrutura(prs)
    s9_pilares(prs)

    # Pilar 1
    pilar_slide(prs, GREEN_MID, "1",
        "Clareza e Identidade Estratégica",
        "O alicerce de tudo.",
        "2 encontros · Mês 1",
        "ENCONTRO 1", "A Virada de Chave",
        ["Diagnóstico de crenças limitantes", "Mapeamento da jornada profissional",
         "Identificação do maior medo", "Definição do estilo de vida desejado",
         "Primeira meta financeira"],
        "ENCONTRO 2", "O Mapa Estratégico",
        ["Definição do nicho e posicionamento", "Frase de posicionamento",
         "Diagnóstico de gargalos", "Limite de atendimentos",
         "3 prioridades de ação imediata"],
    )

    # Pilar 2
    pilar_slide(prs, BROWN_DEEP, "2",
        "Precificação e Estrutura de Negócio",
        "Porque o dinheiro precisa mudar antes da estratégia.",
        "2 encontros · Mês 1-2",
        "ENCONTRO 3", "O Dinheiro que Você Merece",
        ["Cálculo real do custo hora", "Precificação por valor percebido",
         "Pacotes premium", "Política de descontos", "Reserva de segurança"],
        "ENCONTRO 4", "A Estrutura que Sustenta",
        ["PF vs PJ — qual se encaixa", "Pró-labore adequado",
         "Contratos e termos", "Ferramentas de controle financeiro"],
        highlight="✦  Ela sai com o preço novo definido e o compromisso de cobrar na próxima venda."
    )

    # Pilar 3
    pilar_slide(prs, TEAL, "3",
        "Presença Digital e Atração de Pacientes",
        "Atrair pacientes no Instagram é estratégia, narrativa e consistência.",
        "3 encontros · Mês 2-3",
        "ENCONTRO 5", "O Instagram que Converte",
        ["Auditoria do perfil", "Bio com posicionamento", "3 pilares de conteúdo",
         "Estratégia de destaques", "Frequência viável"],
        "ENCONTRO 6", "Conteúdo que Atrai",
        ["Alcance vs conversão", "Roteiro para vídeos curtos",
         "5 tipos de post", "Calendário 30 dias", "Stories com urgência"],
        "ENCONTRO 7", "Do Seguidor ao Paciente",
        ["Fluxo post→DM→WhatsApp", "Script de atendimento",
         "Responder objeções", "Roteiro primeira consulta"],
    )

    # Pilar 4
    pilar_slide(prs, BROWN_MID, "4",
        "Organização e Controle do Negócio",
        "Quando ela já tem o que organizar — a estrutura que sustenta o crescimento.",
        "2 encontros · Mês 3",
        "ENCONTRO 8", "A Estrutura Invisível",
        ["Semana ideal desenhada", "Fluxo do paciente padronizado",
         "Padronização de processos", "Automação no WhatsApp"],
        "ENCONTRO 9", "Gestão e Indicadores",
        ["Os 5 números essenciais", "Carteira ativa vs inativa",
         "Notion/Trello aplicado", "Taxa de ocupação ideal"],
    )

    # Pilar 5
    pilar_slide(prs, GREEN_OLIVE, "5",
        "Consolidação e Próximo Nível",
        "A reta final não é revisão — é construção do futuro.",
        "3 encontros · Mês 4",
        "ENCONTRO 10", "Ajustes e Otimização",
        ["Análise dos indicadores", "Revisão de posicionamento",
         "Ajustes de preço e agenda"],
        "ENCONTRO 11", "Expansão Estratégica",
        ["Consultório próprio vs domiciliar", "Parcerias locais",
         "Construção da marca pessoal", "Planejamento 6 meses"],
        "ENCONTRO 12", "A Fisioterapeuta que Você Virou",
        ["Celebração da jornada", "Carta do Encontro 1",
         "Posicionamento renovado", "Ritual de encerramento"],
    )

    s15_timeline(prs)
    s16_entregaveis(prs)
    s17_transformacao(prs)
    s18_investimento(prs)
    s19_faq(prs)
    s20_cta(prs)

    output = r"c:\Users\trind\Downloads\Claude code\Apresentação Mentoria\Mentoria Plano FT.pptx"
    prs.save(output)
    print(f"PPTX salvo: {output}")
    print(f"Total de slides: {len(prs.slides)}")


if __name__ == "__main__":
    build()
